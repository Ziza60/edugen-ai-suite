"""
build_template.py — Gera o template.pptx para a engine v6 do EduGenAI.
Cada slide usa placeholders {{NOME}} em caixas de texto.
Layout: 16:9 (33.87cm x 19.05cm), paleta navy/gold institucional.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
import copy

# ── Paleta ──────────────────────────────────────────────────────────────────
NAVY       = RGBColor(0x1E, 0x3A, 0x5F)
NAVY_MID   = RGBColor(0x2E, 0x6D, 0xA4)
GOLD       = RGBColor(0xC4, 0x7F, 0x17)
GOLD_LIGHT = RGBColor(0xE8, 0xA0, 0x20)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG   = RGBColor(0xF4, 0xF6, 0xFA)
DARK_BG    = RGBColor(0x0A, 0x16, 0x28)
MUTED      = RGBColor(0x6B, 0x7A, 0x99)
TEXT_DARK  = RGBColor(0x1C, 0x1C, 0x24)

# ── Dimensões slide 16:9 ────────────────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

# Margens
ML = Inches(0.65)
MR = Inches(0.65)
MT = Inches(0.3)
CW = W - ML - MR

CONTENT_Y = Inches(1.35)
FOOTER_Y  = Inches(6.9)
CONTENT_H = FOOTER_Y - CONTENT_Y - Inches(0.1)

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

def blank_layout(prs):
    return prs.slide_layouts[6]  # Blank

def rgb_fill(shape, color: RGBColor):
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, x, y, w, h, color: RGBColor):
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE
    rgb_fill(shape, color)
    shape.line.fill.background()
    return shape

def add_textbox(slide, x, y, w, h, text, font_size, bold=False,
                color=TEXT_DARK, align=PP_ALIGN.LEFT, italic=False,
                font_name="Calibri"):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return txb

def page_chrome(slide, page_label="1 / 45", brand="EduGenAI"):
    """Barra superior navy + footer com paginação e brand."""
    # Barra topo
    add_rect(slide, 0, 0, W, Inches(0.18), NAVY)
    # Acento gold
    add_rect(slide, ML, Inches(0.18), Inches(0.4), Inches(0.04), GOLD)
    # Footer
    add_rect(slide, 0, FOOTER_Y, W, H - FOOTER_Y, NAVY)
    # Paginação (placeholder a ser substituído por string replace)
    add_textbox(slide, W - Inches(1.8), FOOTER_Y + Inches(0.08),
                Inches(1.5), Inches(0.28), page_label,
                8, color=WHITE, align=PP_ALIGN.RIGHT)
    # Brand
    add_textbox(slide, ML, FOOTER_Y + Inches(0.08),
                Inches(3), Inches(0.28), brand,
                8, color=RGBColor(0xA0, 0xB0, 0xCC), align=PP_ALIGN.LEFT)

def chip_label(slide, text):
    """Chip label chip-style: fundo gold arredondado + texto branco."""
    chip_w = max(Inches(1.2), min(Inches(4.2), Pt(len(text) * 7.5)))
    rect = slide.shapes.add_shape(5, ML, Inches(0.22), chip_w, Inches(0.3))  # rounded
    rgb_fill(rect, GOLD)
    rect.line.fill.background()
    tf = rect.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(8)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = "Calibri"

def slide_title(slide, text, y=Inches(0.58), h=Inches(0.75)):
    add_textbox(slide, ML, y, CW, h, text, 24, bold=True,
                color=TEXT_DARK, font_name="Cambria")

# ════════════════════════════════════════════════════════════════════════════
# 0. COVER
# ════════════════════════════════════════════════════════════════════════════
def make_cover(prs):
    sl = prs.slides.add_slide(blank_layout(prs))
    # Fundo escuro
    add_rect(sl, 0, 0, W, H, DARK_BG)
    # Faixa gold inferior
    add_rect(sl, 0, H - Inches(0.6), W, Inches(0.6), GOLD)
    # Barra lateral esquerda navy
    add_rect(sl, 0, 0, Inches(0.5), H, NAVY)
    # Badge
    add_textbox(sl, ML + Inches(0.1), Inches(1.5), Inches(5), Inches(0.35),
                "{{BADGE}}", 9, bold=True, color=GOLD_LIGHT, font_name="Calibri")
    # Título
    add_textbox(sl, ML + Inches(0.1), Inches(2.0), CW - Inches(0.2), Inches(2.5),
                "{{TITLE}}", 40, bold=True, color=WHITE, font_name="Cambria")
    # Tagline
    add_textbox(sl, ML + Inches(0.1), Inches(4.7), CW - Inches(0.2), Inches(0.6),
                "{{TAGLINE}}", 14, color=RGBColor(0xA0, 0xB8, 0xD8), font_name="Calibri")
    # Paginação e brand no footer
    add_rect(sl, 0, H - Inches(0.6), W, Inches(0.6), NAVY)
    add_textbox(sl, W - Inches(1.8), H - Inches(0.52),
                Inches(1.5), Inches(0.3), "1 / 45", 8, color=WHITE, align=PP_ALIGN.RIGHT)
    add_textbox(sl, ML, H - Inches(0.52),
                Inches(3), Inches(0.3), "EduGenAI", 8,
                color=RGBColor(0xA0, 0xB0, 0xCC), align=PP_ALIGN.LEFT)

# ════════════════════════════════════════════════════════════════════════════
# 1. TOC
# ════════════════════════════════════════════════════════════════════════════
def make_toc(prs):
    sl = prs.slides.add_slide(blank_layout(prs))
    add_rect(sl, 0, 0, W, H, LIGHT_BG)
    page_chrome(sl)
    # Painel esquerdo
    panel_w = Inches(3.2)
    add_rect(sl, 0, Inches(0.18), panel_w, H - Inches(0.18) - (H - FOOTER_Y), NAVY)
    add_textbox(sl, Inches(0.2), Inches(0.5), panel_w - Inches(0.3), Inches(0.3),
                "ÍNDICE", 9, bold=True, color=GOLD, font_name="Calibri")
    add_textbox(sl, Inches(0.2), Inches(0.9), panel_w - Inches(0.3), Inches(0.9),
                "Conteúdo\ndo Curso", 20, bold=True, color=WHITE, font_name="Cambria")
    add_textbox(sl, Inches(0.2), Inches(2.0), panel_w - Inches(0.3), Inches(0.4),
                "{{MODULE_COUNT}}", 11, bold=True, color=GOLD_LIGHT, font_name="Calibri")
    # Lista de módulos — 2 colunas
    list_x = panel_w + Inches(0.35)
    list_w = (W - list_x - MR - Inches(0.2)) / 2
    placeholders = [f"{{{{MOD_{i+1}}}}}" for i in range(10)]
    for i, ph in enumerate(placeholders):
        col = i // 5
        row = i % 5
        x = list_x + col * (list_w + Inches(0.2))
        y = CONTENT_Y + row * Inches(0.9)
        # Número do módulo
        add_textbox(sl, x, y, Inches(0.4), Inches(0.4),
                    f"{i+1:02d}", 11, bold=True, color=GOLD, font_name="Calibri")
        # Texto
        add_textbox(sl, x + Inches(0.45), y, list_w - Inches(0.45), Inches(0.8),
                    ph, 11, color=TEXT_DARK, font_name="Calibri")

# ════════════════════════════════════════════════════════════════════════════
# 2. MODULE_COVER
# ════════════════════════════════════════════════════════════════════════════
def make_module_cover(prs):
    sl = prs.slides.add_slide(blank_layout(prs))
    add_rect(sl, 0, 0, W, H, DARK_BG)
    # Barra lateral
    add_rect(sl, 0, 0, Inches(0.5), H, NAVY_MID)
    # Watermark número
    add_textbox(sl, W - Inches(4), Inches(0.1), Inches(3.8), Inches(3.5),
                "{{MODULE_NUMBER}}", 160, bold=True, color=RGBColor(0xFF,0xFF,0xFF),
                align=PP_ALIGN.RIGHT, font_name="Cambria")
    # Label módulo
    add_textbox(sl, ML + Inches(0.1), Inches(1.4), Inches(5), Inches(0.3),
                "{{MODULE_LABEL}}", 9, bold=True, color=GOLD, font_name="Calibri")
    # Título
    add_textbox(sl, ML + Inches(0.1), Inches(1.8), W - ML - Inches(1.3), Inches(1.8),
                "{{TITLE}}", 34, bold=True, color=WHITE, font_name="Cambria")
    # Separador
    add_rect(sl, ML + Inches(0.1), Inches(3.7), Inches(2.2), Inches(0.04), GOLD)
    add_textbox(sl, ML + Inches(0.1), Inches(3.82), Inches(6), Inches(0.22),
                "O QUE VOCÊ VAI APRENDER", 8, bold=True, color=GOLD, font_name="Calibri")
    # Competências
    for i, cy in enumerate([Inches(4.2), Inches(5.1), Inches(6.0)]):
        # bullet
        dot = sl.shapes.add_shape(9, ML + Inches(0.1), cy + Inches(0.08),
                                   Inches(0.12), Inches(0.12))
        rgb_fill(dot, GOLD)
        dot.line.fill.background()
        add_textbox(sl, ML + Inches(0.32), cy, W - ML - Inches(1.5), Inches(0.6),
                    f"{{{{COMP_{i+1}}}}}", 12, color=RGBColor(0xCB, 0xD5, 0xE1),
                    font_name="Calibri")
    # Footer
    add_rect(sl, 0, FOOTER_Y, W, H - FOOTER_Y, NAVY)
    add_textbox(sl, W - Inches(1.8), FOOTER_Y + Inches(0.08),
                Inches(1.5), Inches(0.28), "1 / 45", 8, color=WHITE, align=PP_ALIGN.RIGHT)
    add_textbox(sl, ML, FOOTER_Y + Inches(0.08),
                Inches(3), Inches(0.28), "EduGenAI", 8,
                color=RGBColor(0xA0, 0xB0, 0xCC))

# ════════════════════════════════════════════════════════════════════════════
# 3. BULLETS
# ════════════════════════════════════════════════════════════════════════════
def make_bullets(prs):
    sl = prs.slides.add_slide(blank_layout(prs))
    add_rect(sl, 0, 0, W, H, LIGHT_BG)
    page_chrome(sl)
    chip_label(sl, "{{LABEL}}")
    slide_title(sl, "{{TITLE}}")
    for i in range(5):
        y = CONTENT_Y + i * Inches(0.95)
        add_rect(sl, ML, y + Inches(0.15), Inches(0.08), Inches(0.22), GOLD)
        add_textbox(sl, ML + Inches(0.22), y, CW - Inches(0.22), Inches(0.82),
                    f"{{{{ITEM_{i+1}}}}}", 13, color=TEXT_DARK, font_name="Calibri")

# ════════════════════════════════════════════════════════════════════════════
# 4. CARDS_2
# ════════════════════════════════════════════════════════════════════════════
def make_cards_2(prs):
    sl = prs.slides.add_slide(blank_layout(prs))
    add_rect(sl, 0, 0, W, H, LIGHT_BG)
    page_chrome(sl)
    chip_label(sl, "{{LABEL}}")
    slide_title(sl, "{{TITLE}}")
    gap = Inches(0.25)
    card_w = (CW - gap) / 2
    card_h = Inches(3.8)
    card_y = CONTENT_Y + Inches(0.35)
    colors = [NAVY, NAVY_MID]
    for i in range(2):
        x = ML + i * (card_w + gap)
        add_rect(sl, x, card_y, card_w, card_h, WHITE)
        add_rect(sl, x, card_y, card_w, Inches(0.08), colors[i % 2])
        add_textbox(sl, x + Inches(0.2), card_y + Inches(0.18),
                    card_w - Inches(0.4), Inches(0.55),
                    f"{{{{CARD{i+1}_TITLE}}}}", 14, bold=True,
                    color=colors[i % 2], font_name="Calibri")
        add_textbox(sl, x + Inches(0.2), card_y + Inches(0.82),
                    card_w - Inches(0.4), card_h - Inches(1.1),
                    f"{{{{CARD{i+1}_BODY}}}}", 11, color=TEXT_DARK, font_name="Calibri")

# ════════════════════════════════════════════════════════════════════════════
# 5. CARDS_3
# ════════════════════════════════════════════════════════════════════════════
def make_cards_3(prs):
    sl = prs.slides.add_slide(blank_layout(prs))
    add_rect(sl, 0, 0, W, H, LIGHT_BG)
    page_chrome(sl)
    chip_label(sl, "{{LABEL}}")
    slide_title(sl, "{{TITLE}}")
    gap = Inches(0.2)
    card_w = (CW - gap * 2) / 3
    card_h = Inches(3.6)
    card_y = CONTENT_Y + Inches(0.35)
    colors = [NAVY, NAVY_MID, GOLD]
    for i in range(3):
        x = ML + i * (card_w + gap)
        add_rect(sl, x, card_y, card_w, card_h, WHITE)
        add_rect(sl, x, card_y, card_w, Inches(0.08), colors[i % 3])
        add_textbox(sl, x + Inches(0.15), card_y + Inches(0.18),
                    card_w - Inches(0.3), Inches(0.55),
                    f"{{{{CARD{i+1}_TITLE}}}}", 13, bold=True,
                    color=colors[i % 3], font_name="Calibri")
        add_textbox(sl, x + Inches(0.15), card_y + Inches(0.82),
                    card_w - Inches(0.3), card_h - Inches(1.1),
                    f"{{{{CARD{i+1}_BODY}}}}", 10, color=TEXT_DARK, font_name="Calibri")

# ════════════════════════════════════════════════════════════════════════════
# 6. CARDS_4
# ════════════════════════════════════════════════════════════════════════════
def make_cards_4(prs):
    sl = prs.slides.add_slide(blank_layout(prs))
    add_rect(sl, 0, 0, W, H, LIGHT_BG)
    page_chrome(sl)
    chip_label(sl, "{{LABEL}}")
    slide_title(sl, "{{TITLE}}")
    gap = Inches(0.2)
    card_w = (CW - gap) / 2
    card_h = Inches(1.9)
    colors = [NAVY, NAVY_MID, GOLD, NAVY_MID]
    for i in range(4):
        col = i % 2
        row = i // 2
        x = ML + col * (card_w + gap)
        y = CONTENT_Y + Inches(0.3) + row * (card_h + gap)
        add_rect(sl, x, y, card_w, card_h, WHITE)
        add_rect(sl, x, y, card_w, Inches(0.06), colors[i])
        add_textbox(sl, x + Inches(0.15), y + Inches(0.12),
                    card_w - Inches(0.3), Inches(0.45),
                    f"{{{{CARD{i+1}_TITLE}}}}", 12, bold=True,
                    color=colors[i], font_name="Calibri")
        add_textbox(sl, x + Inches(0.15), y + Inches(0.62),
                    card_w - Inches(0.3), card_h - Inches(0.8),
                    f"{{{{CARD{i+1}_BODY}}}}", 10, color=TEXT_DARK, font_name="Calibri")

# ════════════════════════════════════════════════════════════════════════════
# 7. PROCESS
# ════════════════════════════════════════════════════════════════════════════
def make_process(prs):
    sl = prs.slides.add_slide(blank_layout(prs))
    add_rect(sl, 0, 0, W, H, LIGHT_BG)
    page_chrome(sl)
    chip_label(sl, "{{LABEL}}")
    slide_title(sl, "{{TITLE}}")
    n = 5
    arrow_w = Inches(0.22)
    box_w = (CW - arrow_w * (n - 1)) / n
    box_h = Inches(2.2)
    box_y = CONTENT_Y + (CONTENT_H - box_h) / 2
    for i in range(n):
        x = ML + i * (box_w + arrow_w)
        add_rect(sl, x, box_y, box_w, box_h, WHITE)
        add_rect(sl, x, box_y, box_w, Inches(0.06), NAVY if i % 2 == 0 else NAVY_MID)
        # Badge número
        badge = sl.shapes.add_shape(9, x + Inches(0.15), box_y + Inches(0.15),
                                     Inches(0.38), Inches(0.38))
        rgb_fill(badge, NAVY if i % 2 == 0 else GOLD)
        badge.line.fill.background()
        tf = badge.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = str(i + 1)
        run.font.size = Pt(10)
        run.font.bold = True
        run.font.color.rgb = WHITE
        run.font.name = "Calibri"
        add_textbox(sl, x + Inches(0.1), box_y + Inches(0.65),
                    box_w - Inches(0.2), box_h - Inches(0.8),
                    f"{{{{STEP_{i+1}}}}}", 10, color=TEXT_DARK, font_name="Calibri")
        # Seta
        if i < n - 1:
            ax = x + box_w + Inches(0.04)
            ay = box_y + box_h / 2 - Inches(0.07)
            add_rect(sl, ax, ay, Inches(0.14), Inches(0.04), MUTED)

# ════════════════════════════════════════════════════════════════════════════
# 8. COMPARISON
# ════════════════════════════════════════════════════════════════════════════
def make_comparison(prs):
    sl = prs.slides.add_slide(blank_layout(prs))
    add_rect(sl, 0, 0, W, H, LIGHT_BG)
    page_chrome(sl)
    chip_label(sl, "{{LABEL}}")
    slide_title(sl, "{{TITLE}}")
    col_w = (CW - Inches(0.15)) / 2
    for col, (ph_hdr, phs, color) in enumerate([
        ("LEFT_HEADER",  [f"LEFT_{i+1}"  for i in range(4)], NAVY),
        ("RIGHT_HEADER", [f"RIGHT_{i+1}" for i in range(4)], NAVY_MID),
    ]):
        x = ML + col * (col_w + Inches(0.15))
        add_rect(sl, x, CONTENT_Y, col_w, Inches(0.45), color)
        add_textbox(sl, x + Inches(0.15), CONTENT_Y + Inches(0.08),
                    col_w - Inches(0.3), Inches(0.3),
                    f"{{{{{ph_hdr}}}}}", 13, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, font_name="Calibri")
        for i, ph in enumerate(phs):
            y = CONTENT_Y + Inches(0.55) + i * Inches(1.1)
            add_rect(sl, x, y, col_w, Inches(0.96), WHITE)
            add_textbox(sl, x + Inches(0.15), y + Inches(0.12),
                        col_w - Inches(0.3), Inches(0.75),
                        f"{{{{{ph}}}}}", 11, color=TEXT_DARK, font_name="Calibri")

# ════════════════════════════════════════════════════════════════════════════
# 9. TIMELINE
# ════════════════════════════════════════════════════════════════════════════
def make_timeline(prs):
    sl = prs.slides.add_slide(blank_layout(prs))
    add_rect(sl, 0, 0, W, H, LIGHT_BG)
    page_chrome(sl)
    chip_label(sl, "{{LABEL}}")
    slide_title(sl, "{{TITLE}}")
    # Linha horizontal central
    line_y = CONTENT_Y + CONTENT_H / 2 - Inches(0.02)
    add_rect(sl, ML, line_y, CW, Inches(0.04), NAVY)
    n = 5
    step_w = CW / n
    for i in range(n):
        cx = ML + step_w * i + step_w / 2
        dot = sl.shapes.add_shape(9, cx - Inches(0.18), line_y - Inches(0.16),
                                   Inches(0.36), Inches(0.36))
        rgb_fill(dot, GOLD if i % 2 == 0 else NAVY_MID)
        dot.line.fill.background()
        text_y = line_y - Inches(1.5) if i % 2 == 0 else line_y + Inches(0.35)
        add_textbox(sl, cx - step_w / 2 + Inches(0.1), text_y,
                    step_w - Inches(0.2), Inches(1.1),
                    f"{{{{ITEM_{i+1}}}}}", 10, color=TEXT_DARK,
                    align=PP_ALIGN.CENTER, font_name="Calibri")

# ════════════════════════════════════════════════════════════════════════════
# 10. TWOCOL
# ════════════════════════════════════════════════════════════════════════════
def make_twocol(prs):
    sl = prs.slides.add_slide(blank_layout(prs))
    add_rect(sl, 0, 0, W, H, LIGHT_BG)
    page_chrome(sl)
    chip_label(sl, "{{LABEL}}")
    slide_title(sl, "{{TITLE}}")
    col_w = (CW - Inches(0.5)) / 2
    for col, letters in enumerate(["L", "R"]):
        x = ML + col * (col_w + Inches(0.5))
        add_rect(sl, x, CONTENT_Y, col_w, Inches(0.04), NAVY if col == 0 else NAVY_MID)
        for i in range(4):
            y = CONTENT_Y + Inches(0.18) + i * Inches(1.05)
            add_rect(sl, x, y + Inches(0.2), Inches(0.06), Inches(0.18),
                     GOLD if col == 0 else NAVY_MID)
            add_textbox(sl, x + Inches(0.18), y, col_w - Inches(0.18), Inches(0.88),
                        f"{{{{{letters}{i+1}}}}}", 11, color=TEXT_DARK, font_name="Calibri")

# ════════════════════════════════════════════════════════════════════════════
# 11. TAKEAWAYS
# ════════════════════════════════════════════════════════════════════════════
def make_takeaways(prs):
    sl = prs.slides.add_slide(blank_layout(prs))
    add_rect(sl, 0, 0, W, H, NAVY)
    # Fundo escuro
    add_rect(sl, 0, 0, W, H, DARK_BG)
    # Faixa gold lateral
    add_rect(sl, 0, 0, Inches(0.5), H, GOLD)
    # Chip
    add_textbox(sl, ML + Inches(0.1), Inches(0.3), Inches(5), Inches(0.35),
                "{{LABEL}}", 9, bold=True, color=GOLD, font_name="Calibri")
    # Título
    add_textbox(sl, ML + Inches(0.1), Inches(0.72), CW - Inches(0.2), Inches(0.7),
                "{{TITLE}}", 22, bold=True, color=WHITE, font_name="Cambria")
    for i in range(5):
        y = Inches(1.6) + i * Inches(0.95)
        # Ícone checkmark area
        add_rect(sl, ML + Inches(0.1), y + Inches(0.12),
                 Inches(0.36), Inches(0.36), GOLD)
        add_textbox(sl, ML + Inches(0.6), y, CW - Inches(0.7), Inches(0.82),
                    f"{{{{ITEM_{i+1}}}}}", 12, color=RGBColor(0xE2,0xE8,0xF0),
                    font_name="Calibri")
    # Footer
    add_rect(sl, 0, FOOTER_Y, W, H - FOOTER_Y, NAVY)
    add_textbox(sl, W - Inches(1.8), FOOTER_Y + Inches(0.08),
                Inches(1.5), Inches(0.28), "1 / 45", 8, color=WHITE, align=PP_ALIGN.RIGHT)
    add_textbox(sl, ML, FOOTER_Y + Inches(0.08),
                Inches(3), Inches(0.28), "EduGenAI", 8,
                color=RGBColor(0xA0, 0xB0, 0xCC))

# ════════════════════════════════════════════════════════════════════════════
# 12. CLOSING
# ════════════════════════════════════════════════════════════════════════════
def make_closing(prs):
    sl = prs.slides.add_slide(blank_layout(prs))
    add_rect(sl, 0, 0, W, H, DARK_BG)
    add_rect(sl, 0, 0, Inches(0.5), H, GOLD)
    add_textbox(sl, ML + Inches(0.1), Inches(0.4), CW - Inches(0.2), Inches(0.35),
                "PARABÉNS!", 10, bold=True, color=GOLD, font_name="Calibri")
    add_textbox(sl, ML + Inches(0.1), Inches(0.85), CW - Inches(0.2), Inches(1.2),
                "{{COURSE_TITLE}}", 30, bold=True, color=WHITE, font_name="Cambria")
    add_rect(sl, ML + Inches(0.1), Inches(2.15), Inches(2), Inches(0.04), GOLD)
    add_textbox(sl, ML + Inches(0.1), Inches(2.3), Inches(5), Inches(0.3),
                "PRÓXIMOS PASSOS", 8, bold=True, color=GOLD, font_name="Calibri")
    for i in range(4):
        y = Inches(2.72) + i * Inches(0.92)
        add_rect(sl, ML + Inches(0.1), y + Inches(0.1),
                 Inches(0.35), Inches(0.35), GOLD)
        add_textbox(sl, ML + Inches(0.6), y, CW - Inches(0.7), Inches(0.8),
                    f"{{{{NEXT_{i+1}}}}}", 11, color=RGBColor(0xCB, 0xD5, 0xE1),
                    font_name="Calibri")
    add_rect(sl, 0, FOOTER_Y, W, H - FOOTER_Y, NAVY)
    add_textbox(sl, W - Inches(1.8), FOOTER_Y + Inches(0.08),
                Inches(1.5), Inches(0.28), "1 / 45", 8, color=WHITE, align=PP_ALIGN.RIGHT)
    add_textbox(sl, ML, FOOTER_Y + Inches(0.08),
                Inches(3), Inches(0.28), "EduGenAI", 8,
                color=RGBColor(0xA0, 0xB0, 0xCC))

# ════════════════════════════════════════════════════════════════════════════
# MAIN
# ════════════════════════════════════════════════════════════════════════════

def main():
    prs = new_prs()
    make_cover(prs)          # slide 1 → LAYOUT_INDEX COVER=0
    make_toc(prs)            # slide 2 → TOC=1
    make_module_cover(prs)   # slide 3 → MODULE_COVER=2
    make_bullets(prs)        # slide 4 → BULLETS=3
    make_cards_2(prs)        # slide 5 → CARDS_2=4
    make_cards_3(prs)        # slide 6 → CARDS_3=5
    make_cards_4(prs)        # slide 7 → CARDS_4=6
    make_process(prs)        # slide 8 → PROCESS=7
    make_comparison(prs)     # slide 9 → COMPARISON=8
    make_timeline(prs)       # slide 10 → TIMELINE=9
    make_twocol(prs)         # slide 11 → TWOCOL=10
    make_takeaways(prs)      # slide 12 → TAKEAWAYS=11
    make_closing(prs)        # slide 13 → CLOSING=12

    out = "template.pptx"
    prs.save(out)
    print(f"✅ Template salvo: {out} ({prs.slides.__len__()} slides)")

if __name__ == "__main__":
    main()
